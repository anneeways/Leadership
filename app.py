# Initialize chat history
        if 'chat_history' not in st.session_state:
            st.session_state.chat_history = []
        
        # Add welcome message only once per session
        if 'chatbot_welcomed' not in st.session_state:
            st.session_state.chatbot_welcomed = True
            welcome_msg = f"""
Welcome! I'm your AI consultant for leadership program ROI analysis. I can see you're analyzing a program with {params['participants']} participants and a {import streamlit as st
import pandas as pd
import numpy as np
import json
from datetime import datetime
import plotly.express as px
import plotly.graph_objects as go
from groq import Groq
import os
import io
import base64

# Additional imports for PDF and PowerPoint generation
from reportlab.lib.pagesizes import letter, A4
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, PageBreak
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import inch
from reportlab.lib import colors
from reportlab.graphics.shapes import Drawing
from reportlab.graphics.charts.piecharts import Pie
from reportlab.graphics.charts.barcharts import VerticalBarChart
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Configure Streamlit page
st.set_page_config(
    page_title="Leadership Programme ROI Calculator",
    page_icon="🎯",
    layout="wide",
    initial_sidebar_state="expanded"
)

# Custom CSS for better styling
st.markdown("""
<style>
    .metric-card {
        background: white;
        padding: 1rem;
        border-radius: 0.5rem;
        box-shadow: 0 1px 3px rgba(0,0,0,0.1);
        text-align: center;
    }
    .metric-value {
        font-size: 2rem;
        font-weight: bold;
        margin-bottom: 0.5rem;
    }
    .metric-label {
        color: #666;
        font-weight: 500;
    }
    .metric-status {
        font-size: 0.8rem;
        margin-top: 0.25rem;
    }
    .stTabs [data-baseweb="tab-list"] {
        gap: 2px;
    }
</style>
""", unsafe_allow_html=True)

# Initialize Groq client (optional)
@st.cache_resource
def init_groq():
    try:
        api_key = os.getenv("GROQ_API_KEY") or st.secrets.get("GROQ_API_KEY", "")
        if api_key:
            return Groq(api_key=api_key)
    except:
        pass
    return None

groq_client = init_groq()

# Cost of Capital Templates
COST_OF_CAPITAL_TEMPLATES = {
    'startup': {
        'name': "Startup / High-Growth",
        'rate': 15.0,
        'description': "High risk, high growth expectations"
    },
    'public_large': {
        'name': "Large Public Corporation",
        'rate': 8.0,
        'description': "Established, diversified operations"
    },
    'public_small': {
        'name': "Small/Mid-Cap Public",
        'rate': 10.0,
        'description': "Higher risk than large cap"
    },
    'private_mature': {
        'name': "Mature Private Company",
        'rate': 12.0,
        'description': "Established but less liquid"
    },
    'nonprofit': {
        'name': "Non-Profit Organization",
        'rate': 6.0,
        'description': "Lower risk expectations"
    },
    'government': {
        'name': "Government / Public Sector",
        'rate': 4.0,
        'description': "Risk-free rate plus small premium"
    },
    'private_equity': {
        'name': "Private Equity Backed",
        'rate': 18.0,
        'description': "High return expectations"
    },
    'family_office': {
        'name': "Family Office / UHNW",
        'rate': 7.0,
        'description': "Conservative wealth preservation"
    }
}
INDUSTRY_TEMPLATES = {
    'technology': {
        'name': "Technology",
        'avg_salary': 110000,
        'current_turnover': 22,
        'replacement_cost': 2.0,
        'productivity_gain': 18,
        'retention_improvement': 30,
        'team_performance_gain': 15
    },
    'finance': {
        'name': "Financial Services",
        'avg_salary': 105000,
        'current_turnover': 15,
        'replacement_cost': 1.8,
        'productivity_gain': 12,
        'retention_improvement': 20,
        'team_performance_gain': 10
    },
    'healthcare': {
        'name': "Healthcare",
        'avg_salary': 85000,
        'current_turnover': 20,
        'replacement_cost': 1.6,
        'productivity_gain': 14,
        'retention_improvement': 25,
        'team_performance_gain': 12
    },
    'manufacturing': {
        'name': "Manufacturing",
        'avg_salary': 80000,
        'current_turnover': 16,
        'replacement_cost': 1.4,
        'productivity_gain': 16,
        'retention_improvement': 22,
        'team_performance_gain': 14
    },
    'consulting': {
        'name': "Consulting",
        'avg_salary': 120000,
        'current_turnover': 25,
        'replacement_cost': 2.2,
        'productivity_gain': 20,
        'retention_improvement': 35,
        'team_performance_gain': 18
    }
}

def format_currency(amount):
    """Format amount as currency"""
    return f"${amount:,.0f}"

def format_percentage(value):
    """Format value as percentage"""
    return f"{value:.1f}%"

def get_roi_color_status(roi):
    """Get color and status for ROI"""
    if roi >= 300:
        return "🟢 Excellent", "#10B981"
    elif roi >= 200:
        return "🟡 Good", "#F59E0B"
    elif roi >= 100:
        return "🟠 Moderate", "#EF4444"
    else:
        return "🔴 Review Required", "#DC2626"

def get_payback_color_status(months):
    """Get color and status for payback period"""
    if months <= 12:
        return "🟢 Fast", "#10B981"
    elif months <= 18:
        return "🟡 Moderate", "#F59E0B"
    elif months <= 24:
        return "🟠 Slow", "#EF4444"
    else:
        return "🔴 Very Slow", "#DC2626"

def calculate_roi(params):
    """Calculate ROI and related metrics"""
    # Program Costs
    participant_time_cost = (
        params['participants'] * 
        (params['avg_salary'] * 1.3 / 12) * 
        (params['time_commitment'] / 160) * 
        params['program_duration']
    )
    
    total_program_costs = (
        params['facilitator_costs'] + params['materials_costs'] + 
        params['venue_costs'] + params['travel_costs'] + 
        params['technology_costs'] + params['assessment_costs'] + 
        participant_time_cost
    )
    
    # Annual Benefits
    productivity_benefit = params['participants'] * params['avg_salary'] * (params['productivity_gain'] / 100)
    
    retention_savings = (
        params['participants'] * (params['current_turnover'] / 100) * 
        (params['retention_improvement'] / 100) * params['avg_salary'] * params['replacement_cost']
    )
    
    team_productivity_benefit = (
        params['participants'] * params['team_size'] * 
        (params['avg_salary'] * 0.7) * (params['team_performance_gain'] / 100)
    )
    
    promotion_benefit = (
        params['participants'] * 0.3 * (params['promotion_acceleration'] / 12) * 
        (params['avg_salary'] * 0.2)
    )
    
    decision_benefit = (
        params['participants'] * params['avg_salary'] * 0.1 * 
        (params['decision_quality_gain'] / 100)
    )
    
    total_annual_benefits = (
        productivity_benefit + retention_savings + team_productivity_benefit + 
        promotion_benefit + decision_benefit
    )
    
    # Multi-year analysis
    total_benefits = total_annual_benefits * params['analysis_years']
    net_benefit = total_benefits - total_program_costs
    roi = (net_benefit / total_program_costs) * 100 if total_program_costs > 0 else 0
    payback_months = (total_program_costs / (total_annual_benefits / 12)) if total_annual_benefits > 0 else float('inf')
    
    # NPV calculation with configurable discount rate
    discount_rate = params['discount_rate'] / 100
    npv = -total_program_costs
    
    # Calculate cash flows for each year
    cash_flows = [-total_program_costs]
    for year in range(1, params['analysis_years'] + 1):
        # Adjust for inflation if needed
        inflation_factor = (1 + params['inflation_rate'] / 100) ** year
        adjusted_benefits = total_annual_benefits * inflation_factor
        npv += adjusted_benefits / ((1 + discount_rate) ** year)
        cash_flows.append(adjusted_benefits)
    
    # Calculate IRR (Internal Rate of Return)
    def calculate_irr(cash_flows, max_iter=100, precision=1e-6):
        """Calculate IRR using Newton-Raphson method"""
        if sum(cash_flows) <= 0:
            return None
        
        rate = 0.1  # Initial guess
        for _ in range(max_iter):
            npv_val = sum(cf / (1 + rate) ** i for i, cf in enumerate(cash_flows))
            if abs(npv_val) < precision:
                return rate * 100
            
            # Derivative for Newton-Raphson
            derivative = sum(-i * cf / (1 + rate) ** (i + 1) for i, cf in enumerate(cash_flows))
            if abs(derivative) < precision:
                break
            
            rate = rate - npv_val / derivative
            
            if rate < -0.99:  # Prevent negative rates that are too extreme
                rate = -0.99
        
        return rate * 100 if rate > -1 else None
    
    irr = calculate_irr(cash_flows)
    
    # MIRR (Modified Internal Rate of Return)
    def calculate_mirr(cash_flows, finance_rate, reinvest_rate):
        """Calculate MIRR"""
        n = len(cash_flows) - 1
        
        # Present value of negative cash flows (costs)
        pv_costs = sum(cf / (1 + finance_rate) ** i for i, cf in enumerate(cash_flows) if cf < 0)
        
        # Future value of positive cash flows (benefits)
        fv_benefits = sum(cf * (1 + reinvest_rate) ** (n - i) for i, cf in enumerate(cash_flows) if cf > 0)
        
        if abs(pv_costs) < 1e-6 or fv_benefits < 1e-6:
            return None
            
        mirr = ((fv_benefits / abs(pv_costs)) ** (1/n) - 1) * 100
        return mirr
    
    mirr = calculate_mirr(cash_flows, discount_rate, discount_rate)
    
    # Discounted Payback Period
    def calculate_discounted_payback(cash_flows, discount_rate):
        """Calculate discounted payback period in months"""
        cumulative_pv = 0
        monthly_discount_rate = discount_rate / 12
        
        # Convert annual cash flows to monthly
        monthly_cash_flows = [cash_flows[0]]  # Initial investment
        for i in range(1, len(cash_flows)):
            monthly_benefit = cash_flows[i] / 12
            for month in range(12):
                monthly_cash_flows.append(monthly_benefit)
        
        for month, cash_flow in enumerate(monthly_cash_flows):
            pv_cash_flow = cash_flow / ((1 + monthly_discount_rate) ** month)
            cumulative_pv += pv_cash_flow
            
            if cumulative_pv >= 0 and month > 0:
                return month
        
        return float('inf')  # Never pays back
    
    discounted_payback_months = calculate_discounted_payback(cash_flows, discount_rate)
    
    # Present value of benefits for BCR calculation
    present_value_benefits = sum(cash_flows[i] / (1 + discount_rate) ** i for i in range(1, len(cash_flows)))
    
    benefit_cost_ratio = total_benefits / total_program_costs if total_program_costs > 0 else 0
    
    return {
        'costs': {
            'facilitator': params['facilitator_costs'],
            'materials': params['materials_costs'],
            'venue': params['venue_costs'],
            'travel': params['travel_costs'],
            'technology': params['technology_costs'],
            'assessment': params['assessment_costs'],
            'participant_time': participant_time_cost,
            'total': total_program_costs
        },
        'benefits': {
            'productivity': productivity_benefit,
            'retention': retention_savings,
            'team_performance': team_productivity_benefit,
            'promotion': promotion_benefit,
            'decision_quality': decision_benefit,
            'total_annual': total_annual_benefits,
            'total_multi_year': total_benefits
        },
        'kpis': {
            'roi': roi,
            'payback_months': payback_months,
            'discounted_payback_months': discounted_payback_months,
            'npv': npv,
            'net_benefit': net_benefit,
            'benefit_cost_ratio': benefit_cost_ratio,
            'irr': irr,
            'mirr': mirr,
            'cash_flows': cash_flows
        }
    }

def get_ai_insights(results, params):
    """Get AI-powered insights using Groq"""
    if not groq_client:
        return "AI insights unavailable (Groq API key not configured)"
    
    # List of models to try in order of preference
    models_to_try = [
        "llama3-70b-8192",
        "llama3-8b-8192", 
        "gemma-7b-it",
        "llama3-groq-70b-8192-tool-use-preview"
    ]
    
    prompt = f"""
    Analyze this leadership development program ROI calculation and provide insights:
    
    Program Details:
    - Participants: {params['participants']}
    - Duration: {params['program_duration']} months
    - Investment: {format_currency(results['costs']['total'])}
    
    Key Metrics:
    - ROI: {results['kpis']['roi']:.1f}%
    - Payback: {results['kpis']['payback_months']:.1f} months
    - NPV: {format_currency(results['kpis']['npv'])}
    
    Provide 3-4 bullet points with actionable insights and recommendations.
    Focus on business value, risk mitigation, and optimization opportunities.
    """
    
    for model in models_to_try:
        try:
            response = groq_client.chat.completions.create(
                messages=[{"role": "user", "content": prompt}],
                model=model,
                max_tokens=500
            )
            return response.choices[0].message.content
        except Exception as e:
            if "model_decommissioned" in str(e) or "not found" in str(e):
                continue  # Try next model
            else:
                return f"AI insights unavailable: {str(e)}"
    
    return "AI insights unavailable: All models are currently unavailable"

def get_chatbot_response(user_question, params, results):
    """Get chatbot response using Groq with context about the ROI calculation"""
    if not groq_client:
        return "Chatbot unavailable (Groq API key not configured). Please add your GROQ_API_KEY to use the AI assistant."
    
    # Build comprehensive context
    context = f"""
    You are an expert financial analyst and leadership development consultant helping users understand their ROI calculation for a leadership program. 

    CURRENT PROGRAM SETUP:
    - Participants: {params['participants']}
    - Program Duration: {params['program_duration']} months
    - Average Salary: {format_currency(params['avg_salary'])}
    - Time Commitment: {params['time_commitment']} hours/month
    - Analysis Period: {params['analysis_years']} years
    - Cost of Capital: {params['discount_rate']:.1f}%
    - Expected Inflation: {params['inflation_rate']:.1f}%

    COST BREAKDOWN:
    - Total Investment: {format_currency(results['costs']['total'])}
    - Facilitator Costs: {format_currency(results['costs']['facilitator'])}
    - Materials: {format_currency(results['costs']['materials'])}
    - Venue: {format_currency(results['costs']['venue'])}
    - Travel: {format_currency(results['costs']['travel'])}
    - Technology: {format_currency(results['costs']['technology'])}
    - Assessment: {format_currency(results['costs']['assessment'])}
    - Participant Time: {format_currency(results['costs']['participant_time'])}

    BENEFIT ASSUMPTIONS:
    - Productivity Gain: {params['productivity_gain']:.1f}%
    - Retention Improvement: {params['retention_improvement']:.1f}%
    - Team Performance Gain: {params['team_performance_gain']:.1f}%
    - Decision Quality Gain: {params['decision_quality_gain']:.1f}%

    CURRENT RESULTS:
    - ROI: {results['kpis']['roi']:.1f}%
    - Simple Payback: {results['kpis']['payback_months']:.1f} months
    - Discounted Payback: {results['kpis']['discounted_payback_months']:.1f} months (if not inf)
    - NPV: {format_currency(results['kpis']['npv'])}
    - IRR: {results['kpis']['irr']:.1f}% (if not None)
    - MIRR: {results['kpis']['mirr']:.1f}% (if not None)
    - Benefit-Cost Ratio: {results['kpis']['benefit_cost_ratio']:.1f}:1
    - Annual Benefits: {format_currency(results['benefits']['total_annual'])}

    INSTRUCTIONS:
    - Provide helpful, actionable advice
    - Explain financial concepts in simple terms
    - Suggest specific improvements when asked
    - Use the actual numbers from their calculation
    - Be encouraging but realistic
    - If asked about specific parameters, reference their current settings
    - Help them understand what drives their results
    - Suggest industry benchmarks when relevant
    - Keep responses concise but informative (2-3 paragraphs max)

    User Question: {user_question}
    """
    
    models_to_try = [
        "llama3-70b-8192",
        "llama3-8b-8192", 
        "gemma-7b-it"
    ]
    
    for model in models_to_try:
        try:
            response = groq_client.chat.completions.create(
                messages=[{"role": "user", "content": context}],
                model=model,
                max_tokens=800,
                temperature=0.7
            )
            return response.choices[0].message.content
        except Exception as e:
            if "model_decommissioned" in str(e) or "not found" in str(e):
                continue
            else:
                return f"Chatbot error: {str(e)}"
    
    return "Chatbot temporarily unavailable. All AI models are currently down."

def generate_pdf_report(params, results):
    """Generate PDF report with ROI analysis"""
    buffer = io.BytesIO()
    doc = SimpleDocTemplate(buffer, pagesize=A4)
    styles = getSampleStyleSheet()
    story = []
    
    # Custom styles
    title_style = ParagraphStyle(
        'CustomTitle',
        parent=styles['Heading1'],
        fontSize=24,
        spaceAfter=30,
        textColor=colors.darkblue,
        alignment=1  # Center alignment
    )
    
    heading_style = ParagraphStyle(
        'CustomHeading',
        parent=styles['Heading2'],
        fontSize=16,
        spaceAfter=12,
        textColor=colors.darkblue
    )
    
    # Title page
    story.append(Paragraph("Leadership Programme ROI Analysis", title_style))
    story.append(Spacer(1, 0.5*inch))
    story.append(Paragraph(f"Generated on: {datetime.now().strftime('%B %d, %Y')}", styles['Normal']))
    story.append(Spacer(1, 1*inch))
    
    # Executive Summary
    story.append(Paragraph("Executive Summary", heading_style))
    
    roi_status, _ = get_roi_color_status(results['kpis']['roi'])
    payback_status, _ = get_payback_color_status(results['kpis']['payback_months'])
    
    recommendation = (
        "STRONG BUSINESS CASE - Proceed with implementation" if results['kpis']['roi'] >= 200 else
        "MODERATE BUSINESS CASE - Consider optimization" if results['kpis']['roi'] >= 100 else
        "WEAK BUSINESS CASE - Review assumptions and design"
    )
    
    summary_text = f"""
    <b>Investment:</b> {format_currency(results['costs']['total'])}<br/>
    <b>Annual Benefits:</b> {format_currency(results['benefits']['total_annual'])}<br/>
    <b>Net Benefit:</b> {format_currency(results['kpis']['net_benefit'])}<br/>
    <b>ROI:</b> {results['kpis']['roi']:.0f}% ({roi_status.replace('🟢 ', '').replace('🟡 ', '').replace('🟠 ', '').replace('🔴 ', '')})<br/>
    <b>Payback Period:</b> {results['kpis']['payback_months']:.1f} months<br/>
    <b>Recommendation:</b> {recommendation}
    """
    
    story.append(Paragraph(summary_text, styles['Normal']))
    story.append(Spacer(1, 0.3*inch))
    
    # Program Details
    story.append(Paragraph("Program Details", heading_style))
    
    program_data = [
        ['Metric', 'Value'],
        ['Participants', str(params['participants'])],
        ['Duration', f"{params['program_duration']} months"],
        ['Average Salary', format_currency(params['avg_salary'])],
        ['Time Commitment', f"{params['time_commitment']} hours/month"],
        ['Analysis Period', f"{params['analysis_years']} years"]
    ]
    
    program_table = Table(program_data)
    program_table.setStyle(TableStyle([
        ('BACKGROUND', (0, 0), (-1, 0), colors.grey),
        ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
        ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
        ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
        ('FONTSIZE', (0, 0), (-1, 0), 14),
        ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
        ('BACKGROUND', (0, 1), (-1, -1), colors.beige),
        ('GRID', (0, 0), (-1, -1), 1, colors.black)
    ]))
    
    story.append(program_table)
    story.append(Spacer(1, 0.3*inch))
    
    # Cost Breakdown
    story.append(Paragraph("Cost Breakdown", heading_style))
    
    cost_data = [
        ['Cost Category', 'Amount', 'Percentage'],
        ['Facilitator Costs', format_currency(results['costs']['facilitator']), 
         f"{(results['costs']['facilitator']/results['costs']['total']*100):.1f}%"],
        ['Materials & Content', format_currency(results['costs']['materials']), 
         f"{(results['costs']['materials']/results['costs']['total']*100):.1f}%"],
        ['Venue & Catering', format_currency(results['costs']['venue']), 
         f"{(results['costs']['venue']/results['costs']['total']*100):.1f}%"],
        ['Travel & Accommodation', format_currency(results['costs']['travel']), 
         f"{(results['costs']['travel']/results['costs']['total']*100):.1f}%"],
        ['Technology Platform', format_currency(results['costs']['technology']), 
         f"{(results['costs']['technology']/results['costs']['total']*100):.1f}%"],
        ['Assessment & Evaluation', format_currency(results['costs']['assessment']), 
         f"{(results['costs']['assessment']/results['costs']['total']*100):.1f}%"],
        ['Participant Time', format_currency(results['costs']['participant_time']), 
         f"{(results['costs']['participant_time']/results['costs']['total']*100):.1f}%"],
        ['TOTAL', format_currency(results['costs']['total']), '100.0%']
    ]
    
    cost_table = Table(cost_data)
    cost_table.setStyle(TableStyle([
        ('BACKGROUND', (0, 0), (-1, 0), colors.grey),
        ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
        ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
        ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
        ('FONTNAME', (0, -1), (-1, -1), 'Helvetica-Bold'),
        ('FONTSIZE', (0, 0), (-1, 0), 12),
        ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
        ('BACKGROUND', (0, 1), (-1, -2), colors.beige),
        ('BACKGROUND', (0, -1), (-1, -1), colors.lightgrey),
        ('GRID', (0, 0), (-1, -1), 1, colors.black)
    ]))
    
    story.append(cost_table)
    story.append(PageBreak())
    
    # Benefits Analysis
    story.append(Paragraph("Benefits Analysis", heading_style))
    
    benefit_data = [
        ['Benefit Category', 'Annual Amount', 'Total Value'],
        ['Productivity Improvement', format_currency(results['benefits']['productivity']), 
         format_currency(results['benefits']['productivity'] * params['analysis_years'])],
        ['Retention Savings', format_currency(results['benefits']['retention']), 
         format_currency(results['benefits']['retention'] * params['analysis_years'])],
        ['Team Performance', format_currency(results['benefits']['team_performance']), 
         format_currency(results['benefits']['team_performance'] * params['analysis_years'])],
        ['Promotion Acceleration', format_currency(results['benefits']['promotion']), 
         format_currency(results['benefits']['promotion'] * params['analysis_years'])],
        ['Decision Quality', format_currency(results['benefits']['decision_quality']), 
         format_currency(results['benefits']['decision_quality'] * params['analysis_years'])],
        ['TOTAL', format_currency(results['benefits']['total_annual']), 
         format_currency(results['benefits']['total_multi_year'])]
    ]
    
    benefit_table = Table(benefit_data)
    benefit_table.setStyle(TableStyle([
        ('BACKGROUND', (0, 0), (-1, 0), colors.grey),
        ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
        ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
        ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
        ('FONTNAME', (0, -1), (-1, -1), 'Helvetica-Bold'),
        ('FONTSIZE', (0, 0), (-1, 0), 12),
        ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
        ('BACKGROUND', (0, 1), (-1, -2), colors.beige),
        ('BACKGROUND', (0, -1), (-1, -1), colors.lightgrey),
        ('GRID', (0, 0), (-1, -1), 1, colors.black)
    ]))
    
    story.append(benefit_table)
    story.append(Spacer(1, 0.3*inch))
    
    # Key Performance Indicators
    story.append(Paragraph("Key Performance Indicators", heading_style))
    
    kpi_data = [
        ['KPI', 'Value', 'Status'],
        ['Return on Investment', f"{results['kpis']['roi']:.0f}%", roi_status.replace('🟢 ', '').replace('🟡 ', '').replace('🟠 ', '').replace('🔴 ', '')],
        ['Simple Payback Period', f"{results['kpis']['payback_months']:.1f} months", payback_status.replace('🟢 ', '').replace('🟡 ', '').replace('🟠 ', '').replace('🔴 ', '')],
        ['Discounted Payback Period', f"{results['kpis']['discounted_payback_months']:.1f} months" if results['kpis']['discounted_payback_months'] != float('inf') else 'Never', 'Fast' if results['kpis']['discounted_payback_months'] <= 18 else 'Moderate' if results['kpis']['discounted_payback_months'] <= 30 else 'Slow'],
        ['Net Present Value', format_currency(results['kpis']['npv']), 'Positive' if results['kpis']['npv'] > 0 else 'Negative'],
        ['Benefit-Cost Ratio', f"{results['kpis']['benefit_cost_ratio']:.1f}:1", 'Strong' if results['kpis']['benefit_cost_ratio'] >= 3 else 'Moderate' if results['kpis']['benefit_cost_ratio'] >= 2 else 'Weak'],
        ['Internal Rate of Return', f"{results['kpis']['irr']:.1f}%" if results['kpis']['irr'] else 'N/A', 'Above Hurdle' if results['kpis']['irr'] and results['kpis']['irr'] > params['discount_rate'] else 'Below Hurdle'],
        ['Modified IRR', f"{results['kpis']['mirr']:.1f}%" if results['kpis']['mirr'] else 'N/A', 'Excellent' if results['kpis']['mirr'] and results['kpis']['mirr'] >= 20 else 'Good']
    ]
    
    kpi_table = Table(kpi_data)
    kpi_table.setStyle(TableStyle([
        ('BACKGROUND', (0, 0), (-1, 0), colors.grey),
        ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
        ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
        ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
        ('FONTSIZE', (0, 0), (-1, 0), 12),
        ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
        ('BACKGROUND', (0, 1), (-1, -1), colors.beige),
        ('GRID', (0, 0), (-1, -1), 1, colors.black)
    ]))
    
    story.append(kpi_table)
    story.append(Spacer(1, 0.3*inch))
    
    # Financial Methodology
    story.append(Paragraph("Financial Methodology", heading_style))
    
    methodology_text = f"""
    <b>Key Financial Assumptions:</b><br/>
    • Cost of Capital: {params['discount_rate']:.1f}% (Organization's required return)<br/>
    • Expected Inflation Rate: {params['inflation_rate']:.1f}%<br/>
    • Corporate Tax Rate: {params['tax_rate']:.1f}%<br/>
    • Analysis Period: {params['analysis_years']} years<br/><br/>
    
    <b>Cost of Capital Analysis:</b><br/>
    The {params['discount_rate']:.1f}% cost of capital represents your organization's minimum 
    acceptable return on investment. This rate should reflect your weighted average cost of 
    capital (WACC), considering both debt and equity costs plus appropriate risk premiums.<br/><br/>
    
    <b>Investment Performance vs. Hurdle:</b><br/>
    • Required Return: {params['discount_rate']:.1f}%<br/>
    • Project IRR: {results['kpis']['irr']:.1f}%<br/>
    • Spread: {(results['kpis']['irr'] - params['discount_rate']):+.1f} percentage points<br/><br/>
    
    <b>Calculation Methods:</b><br/>
    • <b>ROI:</b> (Total Benefits - Total Costs) / Total Costs × 100%<br/>
    • <b>NPV:</b> Sum of discounted cash flows minus initial investment<br/>
    • <b>Simple Payback:</b> Time to recover initial investment (nominal terms)<br/>
    • <b>Discounted Payback:</b> Time to recover investment in present value terms<br/>
    • <b>IRR:</b> Discount rate that makes NPV equal to zero<br/>
    • <b>MIRR:</b> Modified IRR assuming reinvestment at cost of capital<br/>
    • <b>Benefit-Cost Ratio:</b> Present value of benefits / Present value of costs<br/><br/>
    
    <b>Benefit Categories:</b><br/>
    • Productivity improvements from enhanced leadership skills<br/>
    • Retention savings from reduced turnover<br/>
    • Team performance gains from better management<br/>
    • Accelerated promotion readiness<br/>
    • Improved decision-making quality<br/>
    """
    
    story.append(Paragraph(methodology_text, styles['Normal']))
    story.append(Spacer(1, 0.3*inch))
    
    # Recommendations
    story.append(Paragraph("Recommendations", heading_style))
    
    recommendations = [
        "• Implement Kirkpatrick Level 3 & 4 evaluation to track behavior change and business results",
        "• Benefits typically materialize 3-6 months post-program completion",
        "• Include manager coaching and 90-day action plans for lasting impact",
        "• World-class leadership programs typically achieve 200-400% ROI",
        "• Start with pilot cohort to validate assumptions before full rollout"
    ]
    
    for rec in recommendations:
        story.append(Paragraph(rec, styles['Normal']))
    
    # Build PDF
    doc.build(story)
    buffer.seek(0)
    return buffer

def generate_powerpoint_report(params, results):
    """Generate PowerPoint presentation with ROI analysis"""
    prs = Presentation()
    
    # Slide 1: Title Slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "Leadership Programme ROI Analysis"
    subtitle.text = f"Business Case & Financial Impact\nGenerated on {datetime.now().strftime('%B %d, %Y')}"
    
    # Slide 2: Executive Summary
    bullet_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    
    title_shape.text = 'Executive Summary'
    
    tf = body_shape.text_frame
    tf.text = f'Investment: {format_currency(results["costs"]["total"])}'
    
    p = tf.add_paragraph()
    p.text = f'Annual Benefits: {format_currency(results["benefits"]["total_annual"])}'
    
    p = tf.add_paragraph()
    p.text = f'ROI: {results["kpis"]["roi"]:.0f}%'
    
    p = tf.add_paragraph()
    p.text = f'Payback Period: {results["kpis"]["payback_months"]:.1f} months'
    
    p = tf.add_paragraph()
    recommendation = (
        "STRONG BUSINESS CASE - Proceed with implementation" if results['kpis']['roi'] >= 200 else
        "MODERATE BUSINESS CASE - Consider optimization" if results['kpis']['roi'] >= 100 else
        "WEAK BUSINESS CASE - Review assumptions and design"
    )
    p.text = f'Recommendation: {recommendation}'
    
    # Slide 3: Program Details
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    
    title_shape.text = 'Program Details'
    
    tf = body_shape.text_frame
    tf.text = f'Participants: {params["participants"]}'
    
    p = tf.add_paragraph()
    p.text = f'Duration: {params["program_duration"]} months'
    
    p = tf.add_paragraph()
    p.text = f'Average Salary: {format_currency(params["avg_salary"])}'
    
    p = tf.add_paragraph()
    p.text = f'Time Commitment: {params["time_commitment"]} hours/month'
    
    p = tf.add_paragraph()
    p.text = f'Analysis Period: {params["analysis_years"]} years'
    
    # Slide 4: Financial Summary
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout
    shapes = slide.shapes
    
    # Add title
    title_shape = shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
    title_frame = title_shape.text_frame
    title_frame.text = "Financial Summary"
    title_frame.paragraphs[0].font.size = Pt(32)
    title_frame.paragraphs[0].font.bold = True
    
    # Add table
    rows = 5
    cols = 2
    left = Inches(2)
    top = Inches(2)
    width = Inches(6)
    height = Inches(3)
    
    table = shapes.add_table(rows, cols, left, top, width, height).table
    
    # Table headers
    table.cell(0, 0).text = 'Metric'
    table.cell(0, 1).text = 'Value'
    
    # Table data
    table.cell(1, 0).text = 'Total Investment'
    table.cell(1, 1).text = format_currency(results['costs']['total'])
    
    table.cell(2, 0).text = 'Annual Benefits'
    table.cell(2, 1).text = format_currency(results['benefits']['total_annual'])
    
    table.cell(3, 0).text = 'Net Benefit'
    table.cell(3, 1).text = format_currency(results['kpis']['net_benefit'])
    
    table.cell(4, 0).text = 'Benefit-Cost Ratio'
    table.cell(4, 1).text = f"{results['kpis']['benefit_cost_ratio']:.1f}:1"
    
    # Slide 5: Key Performance Indicators
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout
    shapes = slide.shapes
    
    # Add title
    title_shape = shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
    title_frame = title_shape.text_frame
    title_frame.text = "Key Performance Indicators"
    title_frame.paragraphs[0].font.size = Pt(32)
    title_frame.paragraphs[0].font.bold = True
    
    # Add KPI table
    rows = 8
    cols = 2
    left = Inches(1)
    top = Inches(1.5)
    width = Inches(8)
    height = Inches(5)
    
    table = shapes.add_table(rows, cols, left, top, width, height).table
    
    # Table headers
    table.cell(0, 0).text = 'Financial Metric'
    table.cell(0, 1).text = 'Value'
    
    # Table data
    table.cell(1, 0).text = 'Return on Investment (ROI)'
    table.cell(1, 1).text = f"{results['kpis']['roi']:.0f}%"
    
    table.cell(2, 0).text = 'Net Present Value (NPV)'
    table.cell(2, 1).text = format_currency(results['kpis']['npv'])
    
    table.cell(3, 0).text = 'Payback Period'
    table.cell(3, 1).text = f"{results['kpis']['payback_months']:.1f} months"
    
    table.cell(4, 0).text = 'Internal Rate of Return (IRR)'
    table.cell(4, 1).text = f"{results['kpis']['irr']:.1f}%" if results['kpis']['irr'] else 'N/A'
    
    table.cell(5, 0).text = 'Modified IRR (MIRR)'
    table.cell(5, 1).text = f"{results['kpis']['mirr']:.1f}%" if results['kpis']['mirr'] else 'N/A'
    
    table.cell(6, 0).text = 'Benefit-Cost Ratio'
    table.cell(6, 1).text = f"{results['kpis']['benefit_cost_ratio']:.1f}:1"
    
    table.cell(7, 0).text = 'Profitability Index'
    table.cell(7, 1).text = f"{results['kpis']['profitability_index']:.2f}"
    
    # Slide 6: Financial Methodology
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    
    title_shape.text = 'Financial Methodology'
    
    tf = body_shape.text_frame
    tf.text = f'Cost of Capital: {params["discount_rate"]:.1f}% (Organization hurdle rate)'
    
    p = tf.add_paragraph()
    p.text = f'Analysis Period: {params["analysis_years"]} years with {params["inflation_rate"]:.1f}% inflation'
    
    p = tf.add_paragraph()
    irr_vs_hurdle = results['kpis']['irr'] - params['discount_rate'] if results['kpis']['irr'] else 0
    p.text = f'IRR vs. Cost of Capital: {irr_vs_hurdle:+.1f} percentage points'
    
    p = tf.add_paragraph()
    p.text = 'NPV calculated using discounted cash flows'
    
    p = tf.add_paragraph()
    p.text = 'IRR assumes reinvestment at hurdle rate'
    
    p = tf.add_paragraph()
    p.text = 'Benefits include productivity, retention, and team gains'
    
    # Slide 7: Recommendations
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    
    title_shape.text = 'Recommendations'
    
    tf = body_shape.text_frame
    tf.text = 'Implement Kirkpatrick Level 3 & 4 evaluation'
    
    recommendations = [
        'Benefits typically materialize 3-6 months post-program',
        'Include manager coaching and 90-day action plans',
        'World-class programs achieve 200-400% ROI',
        'Start with pilot cohort to validate assumptions'
    ]
    
    for rec in recommendations:
        p = tf.add_paragraph()
        p.text = rec
    
    # Save to buffer
    buffer = io.BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    return buffer

def main():
    # Header
    st.markdown("""
    <div style='background: linear-gradient(135deg, #667eea 0%, #764ba2 100%); 
                padding: 2rem; border-radius: 10px; margin-bottom: 2rem;'>
        <h1 style='color: white; margin: 0; font-size: 2.5rem;'>🎯 Leadership Programme ROI Calculator</h1>
        <p style='color: rgba(255,255,255,0.8); margin: 0.5rem 0 0 0; font-size: 1.2rem;'>
            Calculate ROI, payback period, and build compelling business cases for leadership development initiatives
        </p>
    </div>
    """, unsafe_allow_html=True)
    
    # Initialize session state
    if 'params' not in st.session_state:
        st.session_state.params = {}
    
    # Ensure all parameters exist with defaults
    default_params = {
        # Program Parameters
        'participants': 20,
        'program_duration': 6,
        'avg_salary': 95000,
        'time_commitment': 15,
        'analysis_years': 3,
        
        # Program Costs
        'facilitator_costs': 75000,
        'materials_costs': 15000,
        'venue_costs': 25000,
        'travel_costs': 30000,
        'technology_costs': 12000,
        'assessment_costs': 8000,
        
        # Benefit Assumptions
        'productivity_gain': 15,
        'retention_improvement': 25,
        'promotion_acceleration': 6,
        'team_performance_gain': 12,
        'decision_quality_gain': 20,
        
        # Industry Benchmarks
        'current_turnover': 18,
        'replacement_cost': 1.5,
        'team_size': 8,
        
        # Financial Parameters
        'discount_rate': 8.0,
        'tax_rate': 25.0,
        'inflation_rate': 3.0
    }
    
    # Update session state with any missing parameters
    for key, value in default_params.items():
        if key not in st.session_state.params:
            st.session_state.params[key] = value
    
    # Sidebar for inputs
    with st.sidebar:
        st.header("📊 Program Configuration")
        
        # Industry Templates
        st.subheader("🏭 Industry Templates")
        industry_options = [''] + [template['name'] for template in INDUSTRY_TEMPLATES.values()]
        selected_industry = st.selectbox("Select Industry Template", industry_options)
        
        if selected_industry:
            # Find the template key
            template_key = None
            for key, template in INDUSTRY_TEMPLATES.items():
                if template['name'] == selected_industry:
                    template_key = key
                    break
            
            if template_key and st.button("Apply Template"):
                template = INDUSTRY_TEMPLATES[template_key]
                st.session_state.params.update({
                    'avg_salary': template['avg_salary'],
                    'current_turnover': template['current_turnover'],
                    'replacement_cost': template['replacement_cost'],
                    'productivity_gain': template['productivity_gain'],
                    'retention_improvement': template['retention_improvement'],
                    'team_performance_gain': template['team_performance_gain']
                })
                st.success(f"Applied {selected_industry} template!")
        
        st.divider()
        
        # Program Parameters
        st.subheader("📈 Program Parameters")
        st.session_state.params['participants'] = st.number_input(
            "Participants", min_value=1, value=st.session_state.params['participants']
        )
        st.session_state.params['program_duration'] = st.number_input(
            "Duration (months)", min_value=1, value=st.session_state.params['program_duration']
        )
        st.session_state.params['avg_salary'] = st.number_input(
            "Average Salary ($)", min_value=0, value=st.session_state.params['avg_salary'], step=5000
        )
        st.session_state.params['time_commitment'] = st.number_input(
            "Time Commitment (hours/month)", min_value=1, value=st.session_state.params['time_commitment']
        )
        st.session_state.params['analysis_years'] = st.number_input(
            "Analysis Period (years)", min_value=1, max_value=10, value=st.session_state.params['analysis_years']
        )
        
        st.divider()
        
        # Program Costs
        st.subheader("💰 Program Costs")
        cost_fields = [
            ('facilitator_costs', 'Facilitator Costs ($)', 5000),
            ('materials_costs', 'Materials & Content ($)', 1000),
            ('venue_costs', 'Venue & Catering ($)', 1000),
            ('travel_costs', 'Travel & Accommodation ($)', 1000),
            ('technology_costs', 'Technology Platform ($)', 1000),
            ('assessment_costs', 'Assessment & Evaluation ($)', 1000)
        ]
        
        for field, label, step in cost_fields:
            st.session_state.params[field] = st.number_input(
                label, min_value=0, value=st.session_state.params[field], step=step
            )
        
        st.divider()
        
        # Cost of Capital Selection
        st.subheader("💰 Cost of Capital")
        
        # Cost of Capital Templates
        cost_of_capital_options = ['Custom'] + [template['name'] for template in COST_OF_CAPITAL_TEMPLATES.values()]
        selected_cost_template = st.selectbox(
            "Organization Type", 
            cost_of_capital_options,
            help="Select your organization type for appropriate cost of capital"
        )
        
        if selected_cost_template != 'Custom':
            # Find the template key
            template_key = None
            for key, template in COST_OF_CAPITAL_TEMPLATES.items():
                if template['name'] == selected_cost_template:
                    template_key = key
                    break
            
            if template_key:
                template = COST_OF_CAPITAL_TEMPLATES[template_key]
                st.info(f"**{template['name']}**: {template['description']}")
                
                if st.button("Apply Cost of Capital"):
                    st.session_state.params['discount_rate'] = template['rate']
                    st.success(f"Applied {template['rate']:.1f}% cost of capital!")
                    st.rerun()
        
        # Manual cost of capital input
        st.session_state.params['discount_rate'] = st.number_input(
            "Cost of Capital / Discount Rate (%)", 
            min_value=0.0, max_value=25.0, 
            value=st.session_state.params['discount_rate'], 
            step=0.5,
            help="Your organization's weighted average cost of capital (WACC)"
        )
        
        # Cost of capital guidance
        with st.expander("💡 Cost of Capital Guidance"):
            st.markdown("""
            **What is Cost of Capital?**
            The minimum return your organization expects from investments, representing:
            - **Cost of Debt**: Interest rates on borrowed money
            - **Cost of Equity**: Expected returns for shareholders
            - **Risk Premium**: Additional return for business risk
            
            **Typical Ranges by Organization:**
            - **Government/Public Sector**: 3-5%
            - **Large Public Companies**: 6-10%
            - **Small/Medium Private**: 10-15%
            - **Startups/High Growth**: 15-20%
            - **Private Equity**: 18-25%
            
            **Consider Your:**
            - Industry risk level
            - Company size and maturity
            - Financial leverage
            - Market conditions
            - Investment time horizon
            """)
        
        st.divider()
        
        # Financial Parameters
        st.subheader("🏦 Other Financial Parameters")
        st.session_state.params['tax_rate'] = st.number_input(
            "Corporate Tax Rate (%)", 
            min_value=0.0, max_value=50.0, 
            value=st.session_state.params['tax_rate'], 
            step=1.0,
            help="Corporate tax rate for tax shield calculations"
        )
        st.session_state.params['inflation_rate'] = st.number_input(
            "Expected Inflation Rate (%)", 
            min_value=0.0, max_value=10.0, 
            value=st.session_state.params['inflation_rate'], 
            step=0.5,
            help="Expected annual inflation rate for benefit adjustments"
        )
    
    # Calculate results
    results = calculate_roi(st.session_state.params)
    
    # Main content tabs
    tab1, tab2, tab3, tab4, tab5 = st.tabs(["📊 Dashboard", "⚙️ Assumptions", "📈 Analysis", "🤖 AI Insights", "💬 AI Assistant"])
    
    with tab1:
        # Cost of Capital Indicator
        current_rate = st.session_state.params['discount_rate']
        
        # Find if current rate matches any template
        matching_template = None
        for key, template in COST_OF_CAPITAL_TEMPLATES.items():
            if abs(template['rate'] - current_rate) < 0.1:
                matching_template = template
                break
        
        if matching_template:
            st.info(f"💰 **Cost of Capital**: {current_rate:.1f}% ({matching_template['name']}) - {matching_template['description']}")
        else:
            st.info(f"💰 **Cost of Capital**: {current_rate:.1f}% (Custom rate for your organization)")
        
        # Key Metrics Dashboard
        col1, col2 = st.columns([3, 1])
        with col1:
            st.subheader("🎯 Key Performance Indicators")
        with col2:
            if st.button("💬 Ask AI about KPIs", key="help_kpis"):
                question = "Can you help me understand what these KPI results mean for my business case?"
                if 'chat_history' not in st.session_state:
                    st.session_state.chat_history = []
                st.session_state.chat_history.append(("user", question))
                response = get_chatbot_response(question, st.session_state.params, results)
                st.session_state.chat_history.append(("assistant", response))
                st.success("💬 Question added to AI Assistant! Check the 'AI Assistant' tab for the response.")
        
        col1, col2, col3, col4 = st.columns(4)
        
        with col1:
            roi_status, roi_color = get_roi_color_status(results['kpis']['roi'])
            st.metric(
                "Return on Investment",
                f"{results['kpis']['roi']:.0f}%",
                delta=roi_status
            )
        
        with col2:
            payback_status, payback_color = get_payback_color_status(results['kpis']['payback_months'])
            st.metric(
                "Payback Period",
                f"{results['kpis']['payback_months']:.1f} months",
                delta=payback_status
            )
        
        with col3:
            npv_status = "🟢 Positive" if results['kpis']['npv'] > 0 else "🔴 Negative"
            st.metric(
                "Net Present Value",
                format_currency(results['kpis']['npv']),
                delta=npv_status
            )
        
        with col4:
            bcr = results['kpis']['benefit_cost_ratio']
            bcr_status = "🟢 Strong" if bcr >= 3 else "🟡 Moderate" if bcr >= 2 else "🔴 Weak"
            st.metric(
                "Benefit-Cost Ratio",
                f"{bcr:.1f}:1",
                delta=bcr_status
            )
        
        # Additional Financial Metrics
        st.subheader("📈 Advanced Financial Metrics")
        
        col1, col2, col3, col4 = st.columns(4)
        
        with col1:
            irr_value = results['kpis']['irr']
            if irr_value is not None:
                irr_status = "🟢 Excellent" if irr_value >= 25 else "🟡 Good" if irr_value >= 15 else "🟠 Moderate"
                st.metric(
                    "Internal Rate of Return",
                    f"{irr_value:.1f}%",
                    delta=irr_status
                )
            else:
                st.metric("Internal Rate of Return", "N/A", delta="🔴 Cannot calculate")
        
        with col2:
            mirr_value = results['kpis']['mirr']
            if mirr_value is not None:
                mirr_status = "🟢 Excellent" if mirr_value >= 20 else "🟡 Good" if mirr_value >= 12 else "🟠 Moderate"
                st.metric(
                    "Modified IRR",
                    f"{mirr_value:.1f}%",
                    delta=mirr_status
                )
            else:
                st.metric("Modified IRR", "N/A", delta="🔴 Cannot calculate")
        
        with col3:
            dpb_months = results['kpis']['discounted_payback_months']
            if dpb_months != float('inf'):
                dpb_status = "🟢 Fast" if dpb_months <= 18 else "🟡 Moderate" if dpb_months <= 30 else "🟠 Slow"
                st.metric(
                    "Discounted Payback",
                    f"{dpb_months:.1f} months",
                    delta=dpb_status
                )
            else:
                st.metric("Discounted Payback", "Never", delta="🔴 No payback")
        
        with col4:
            discount_rate = st.session_state.params['discount_rate']
            irr_value = results['kpis']['irr']
            
            if irr_value is not None:
                spread = irr_value - discount_rate
                hurdle_comparison = "🟢 Above hurdle" if spread > 0 else "🔴 Below hurdle"
                st.metric(
                    "vs. Cost of Capital",
                    f"+{spread:.1f}%" if spread > 0 else f"{spread:.1f}%",
                    delta=f"Hurdle: {discount_rate:.1f}%"
                )
            else:
                st.metric(
                    "Cost of Capital",
                    f"{discount_rate:.1f}%",
                    delta="Minimum required return"
                )
        
        # Methodology Expander
        with st.expander("📚 Calculation Methodology", expanded=False):
            st.markdown("""
            ### Financial Metrics Explained
            
            **Return on Investment (ROI)**
            ```
            ROI = (Total Benefits - Total Costs) / Total Costs × 100%
            ```
            Measures the efficiency of the investment relative to its cost.
            
            **Net Present Value (NPV)**
            ```
            NPV = Σ(Cash Flow_t / (1 + r)^t) - Initial Investment
            ```
            Where r = discount rate, t = time period. Accounts for time value of money.
            
            **Internal Rate of Return (IRR)**
            ```
            0 = Σ(Cash Flow_t / (1 + IRR)^t)
            ```
            The discount rate that makes NPV equal to zero.
            
            **Modified Internal Rate of Return (MIRR)**
            ```
            MIRR = (Future Value of Positive Flows / Present Value of Negative Flows)^(1/n) - 1
            ```
            More realistic than IRR as it assumes reinvestment at the cost of capital.
            
            **Payback Period (Simple)**
            ```
            Payback = Initial Investment / Annual Cash Flow
            ```
            Time required to recover the initial investment (ignores time value of money).
            
            **Discounted Payback Period**
            ```
            Cumulative PV of Cash Flows ≥ 0
            ```
            Time to recover investment using present value of cash flows (accounts for time value).
            
            **Benefit-Cost Ratio (BCR)**
            ```
            BCR = Present Value of Benefits / Present Value of Costs
            ```
            Values > 1.0 indicate positive value creation.
            
            ### Key Assumptions:
            - **Cost of Capital:** {:.1f}% (Your organization's required return)
            - **Inflation Rate:** {:.1f}% (Applied to future benefits)
            - **Tax Rate:** {:.1f}% (For tax considerations)
            - **Analysis Period:** {} years
            - **Participant Time Cost:** Loaded salary rate (salary × 1.3) to include benefits
            
            ### Cost of Capital Context:
            Current rate of {:.1f}% represents your organization's minimum acceptable return.
            The IRR of {:.1f}% {} this hurdle rate by {:.1f} percentage points.
            """.format(
                st.session_state.params['discount_rate'],
                st.session_state.params['inflation_rate'], 
                st.session_state.params['tax_rate'],
                st.session_state.params['analysis_years'],
                st.session_state.params['discount_rate'],
                results['kpis']['irr'] if results['kpis']['irr'] else 0,
                "exceeds" if results['kpis']['irr'] and results['kpis']['irr'] > st.session_state.params['discount_rate'] else "falls short of",
                abs((results['kpis']['irr'] or 0) - st.session_state.params['discount_rate'])
            ))
        
        st.divider()
        
        # Cost and Benefit Breakdown
        col1, col2 = st.columns(2)
        
        with col1:
            col1a, col1b = st.columns([3, 1])
            with col1a:
                st.subheader("💰 Cost Breakdown")
            with col1b:
                if st.button("💬 Optimize Costs", key="help_costs", help="Ask AI how to optimize costs"):
                    question = "Looking at my cost breakdown, how can I optimize my investment while maintaining program quality?"
                    if 'chat_history' not in st.session_state:
                        st.session_state.chat_history = []
                    st.session_state.chat_history.append(("user", question))
                    response = get_chatbot_response(question, st.session_state.params, results)
                    st.session_state.chat_history.append(("assistant", response))
                    st.success("💬 Question sent to AI Assistant!")
            cost_data = {
                'Category': ['Facilitator', 'Materials', 'Venue', 'Travel', 'Technology', 'Assessment', 'Participant Time'],
                'Amount': [
                    results['costs']['facilitator'],
                    results['costs']['materials'],
                    results['costs']['venue'],
                    results['costs']['travel'],
                    results['costs']['technology'],
                    results['costs']['assessment'],
                    results['costs']['participant_time']
                ]
            }
            
            fig_costs = px.pie(
                values=cost_data['Amount'],
                names=cost_data['Category'],
                title=f"Total Investment: {format_currency(results['costs']['total'])}"
            )
            st.plotly_chart(fig_costs, use_container_width=True)
        
        with col2:
            col2a, col2b = st.columns([3, 1])
            with col2a:
                st.subheader("📈 Benefit Breakdown")
            with col2b:
                if st.button("💬 Validate Benefits", key="help_benefits", help="Ask AI about benefit assumptions"):
                    question = "Are my benefit assumptions realistic? How can I strengthen the business case for these expected returns?"
                    if 'chat_history' not in st.session_state:
                        st.session_state.chat_history = []
                    st.session_state.chat_history.append(("user", question))
                    response = get_chatbot_response(question, st.session_state.params, results)
                    st.session_state.chat_history.append(("assistant", response))
                    st.success("💬 Question sent to AI Assistant!")
            benefit_data = {
                'Category': ['Productivity', 'Retention', 'Team Performance', 'Promotions', 'Decision Quality'],
                'Annual Amount': [
                    results['benefits']['productivity'],
                    results['benefits']['retention'],
                    results['benefits']['team_performance'],
                    results['benefits']['promotion'],
                    results['benefits']['decision_quality']
                ]
            }
            
            fig_benefits = px.bar(
                x=benefit_data['Category'],
                y=benefit_data['Annual Amount'],
                title=f"Annual Benefits: {format_currency(results['benefits']['total_annual'])}"
            )
            fig_benefits.update_layout(xaxis_tickangle=-45)
            st.plotly_chart(fig_benefits, use_container_width=True)
    
    with tab2:
        st.subheader("⚙️ Impact Assumptions")
        
        col1, col2 = st.columns(2)
        
        with col1:
            st.session_state.params['productivity_gain'] = st.slider(
                "Productivity Improvement (%)", 0, 50, st.session_state.params['productivity_gain'],
                help="Industry average: 10-20%"
            )
            st.session_state.params['retention_improvement'] = st.slider(
                "Retention Improvement (%)", 0, 50, st.session_state.params['retention_improvement'],
                help="Industry average: 15-30%"
            )
            st.session_state.params['team_performance_gain'] = st.slider(
                "Team Performance Gain (%)", 0, 30, st.session_state.params['team_performance_gain'],
                help="Industry average: 8-15%"
            )
        
        with col2:
            st.session_state.params['decision_quality_gain'] = st.slider(
                "Decision Quality Improvement (%)", 0, 40, st.session_state.params['decision_quality_gain'],
                help="Industry average: 15-25%"
            )
            st.session_state.params['current_turnover'] = st.slider(
                "Current Turnover Rate (%)", 0, 50, st.session_state.params['current_turnover'],
                help="Annual turnover rate"
            )
            st.session_state.params['replacement_cost'] = st.slider(
                "Replacement Cost Multiple", 0.5, 3.0, st.session_state.params['replacement_cost'], 0.1,
                help="Multiple of salary to replace employee"
            )
    
    with tab3:
        st.subheader("📊 Scenario Analysis")
        
        # Scenario comparison
        scenarios = {
            'Conservative (-25%)': 0.75,
            'Realistic (Current)': 1.0,
            'Optimistic (+25%)': 1.25
        }
        
        scenario_data = []
        for scenario_name, multiplier in scenarios.items():
            # Temporarily calculate with different scenario
            temp_params = st.session_state.params.copy()
            # Scale the benefit assumptions
            temp_params['productivity_gain'] = temp_params['productivity_gain'] * multiplier
            temp_params['retention_improvement'] = temp_params['retention_improvement'] * multiplier
            temp_params['team_performance_gain'] = temp_params['team_performance_gain'] * multiplier
            temp_params['decision_quality_gain'] = temp_params['decision_quality_gain'] * multiplier
            temp_results = calculate_roi(temp_params)
            
            dpb = temp_results['kpis']['discounted_payback_months']
            dpb_text = f"{dpb:.1f}" if dpb != float('inf') else "Never"
            
            scenario_data.append({
                'Scenario': scenario_name,
                'ROI (%)': f"{temp_results['kpis']['roi']:.0f}%",
                'Simple Payback (months)': f"{temp_results['kpis']['payback_months']:.1f}",
                'Discounted Payback (months)': dpb_text,
                'Net Benefit': format_currency(temp_results['kpis']['net_benefit'])
            })
        
        df_scenarios = pd.DataFrame(scenario_data)
        st.table(df_scenarios)
        
        st.divider()
        
        # Cash flow visualization
        st.subheader("💰 Cash Flow Analysis")
        
        years = list(range(0, st.session_state.params['analysis_years'] + 1))
        cumulative_cash_flow = [-results['costs']['total']]
        
        for year in range(1, st.session_state.params['analysis_years'] + 1):
            cumulative_cash_flow.append(
                cumulative_cash_flow[-1] + results['benefits']['total_annual']
            )
        
        fig_cashflow = go.Figure()
        fig_cashflow.add_trace(go.Scatter(
            x=years,
            y=cumulative_cash_flow,
            mode='lines+markers',
            name='Cumulative Cash Flow',
            line=dict(width=3)
        ))
        fig_cashflow.add_hline(y=0, line_dash="dash", line_color="red", 
                              annotation_text="Break-even")
        fig_cashflow.update_layout(
            title="Cumulative Cash Flow Over Time",
            xaxis_title="Year",
            yaxis_title="Cumulative Cash Flow ($)",
            hovermode='x'
        )
        st.plotly_chart(fig_cashflow, use_container_width=True)
        
        st.divider()
        
        # Discount Rate Sensitivity Analysis
        st.subheader("📊 Discount Rate Sensitivity Analysis")
        
        # Generate sensitivity data
        discount_rates = [2, 4, 6, 8, 10, 12, 15, 20]
        sensitivity_data = []
        
        for rate in discount_rates:
            # Temporarily calculate with different discount rate
            temp_params = st.session_state.params.copy()
            temp_params['discount_rate'] = rate
            temp_results = calculate_roi(temp_params)
            
            sensitivity_data.append({
                'Discount Rate (%)': rate,
                'NPV': temp_results['kpis']['npv'],
                'ROI (%)': temp_results['kpis']['roi'],
                'BCR': temp_results['kpis']['benefit_cost_ratio']
            })
        
        df_sensitivity = pd.DataFrame(sensitivity_data)
        
        # Create sensitivity chart
        fig_sensitivity = go.Figure()
        
        fig_sensitivity.add_trace(go.Scatter(
            x=df_sensitivity['Discount Rate (%)'],
            y=df_sensitivity['NPV'],
            mode='lines+markers',
            name='NPV ($)',
            yaxis='y',
            line=dict(color='blue', width=3)
        ))
        
        # Add current discount rate line
        current_rate = st.session_state.params['discount_rate']
        fig_sensitivity.add_vline(
            x=current_rate, 
            line_dash="dash", 
            line_color="red",
            annotation_text=f"Current Rate: {current_rate}%"
        )
        
        fig_sensitivity.add_hline(y=0, line_dash="dash", line_color="gray", annotation_text="Break-even NPV")
        
        fig_sensitivity.update_layout(
            title="NPV Sensitivity to Discount Rate Changes",
            xaxis_title="Discount Rate (%)",
            yaxis_title="Net Present Value ($)",
            hovermode='x unified'
        )
        
        st.plotly_chart(fig_sensitivity, use_container_width=True)
        
        # Sensitivity table
        st.markdown("**Sensitivity Analysis Table**")
        
        # Highlight current rate row
        def highlight_current_rate(row):
            if abs(row['Discount Rate (%)'] - current_rate) < 0.1:
                return ['background-color: lightblue'] * len(row)
            return [''] * len(row)
        
        # Format the dataframe for display
        df_display = df_sensitivity.copy()
        df_display['NPV'] = df_display['NPV'].apply(lambda x: format_currency(x))
        df_display['ROI (%)'] = df_display['ROI (%)'].apply(lambda x: f"{x:.1f}%")
        df_display['BCR'] = df_display['BCR'].apply(lambda x: f"{x:.2f}")
        
        st.dataframe(df_display.style.apply(highlight_current_rate, axis=1), use_container_width=True)
    
    with tab4:
        st.subheader("🤖 AI-Powered Insights")
        
        if st.button("Generate AI Insights", type="primary"):
            with st.spinner("Analyzing your business case..."):
                insights = get_ai_insights(results, st.session_state.params)
                st.markdown(insights)
        
        st.divider()
        
        st.subheader("💡 Best Practices & Recommendations")
        
        recommendations = [
            "**Measurement:** Implement Kirkpatrick Level 3 & 4 evaluation to track behavior change and business results",
            "**Timeline:** Benefits typically materialize 3-6 months post-program completion",
            "**Sustainability:** Include manager coaching and 90-day action plans for lasting impact",
            "**Benchmark:** World-class leadership programs typically achieve 200-400% ROI",
            "**Risk Mitigation:** Start with pilot cohort to validate assumptions before full rollout"
        ]
        
        for rec in recommendations:
            st.markdown(f"✓ {rec}")
    
    with tab5:
        st.subheader("💬 AI Assistant - Your Personal ROI Consultant")
        
        # Initialize chat history
        if 'chat_history' not in st.session_state:
            st.session_state.chat_history = []
            # Add welcome message
            welcome_msg = f"""
            Welcome! I'm your AI consultant for leadership program ROI analysis. I can see you're analyzing a program with {params['participants']} participants and a {format_currency(results['costs']['total'])} investment.

            Your current ROI is {results['kpis']['roi']:.1f}% with a {results['kpis']['payback_months']:.1f}-month payback period. 

            I can help you:
            • 📊 Interpret your financial results
            • 💡 Suggest improvements to your business case  
            • ⚖️ Compare against industry benchmarks
            • 🎯 Optimize costs and assumptions
            • 📋 Prepare for stakeholder presentations

            What would you like to explore first?
            """
            st.session_state.chat_history.append(("assistant", welcome_msg))
        
        # Quick action buttons
        st.markdown("**🚀 Quick Help:**")
        col1, col2, col3, col4 = st.columns(4)
        
        with col1:
            if st.button("💡 Explain My Results", key="explain_results"):
                question = "Can you explain what my current ROI results mean and whether this is a good investment?"
                st.session_state.pending_question = question
        
        with col2:
            if st.button("📈 Improve My ROI", key="improve_roi"):
                question = "How can I improve my ROI? What parameters should I focus on?"
                st.session_state.pending_question = question
        
        with col3:
            if st.button("⚖️ Industry Benchmark", key="benchmark"):
                question = "How do my results compare to industry benchmarks for leadership programs?"
                st.session_state.pending_question = question
        
        with col4:
            if st.button("🎯 Optimize Costs", key="optimize_costs"):
                question = "Which cost categories should I focus on to optimize my investment?"
                st.session_state.pending_question = question
        
        st.divider()
        
        # Chat interface
        st.markdown("**💬 Ask me anything about your ROI analysis:**")
        
        # Display chat history
        for i, (role, message) in enumerate(st.session_state.chat_history):
            if role == "user":
                st.markdown(f"**🙋 You:** {message}")
            else:
                st.markdown(f"**🤖 Assistant:** {message}")
            
            if i < len(st.session_state.chat_history) - 1:
                st.markdown("---")
        
        # Handle pending question from quick buttons
        if hasattr(st.session_state, 'pending_question'):
            question = st.session_state.pending_question
            delattr(st.session_state, 'pending_question')
            
            # Add user question to history
            st.session_state.chat_history.append(("user", question))
            
            # Get AI response
            with st.spinner("🤖 Analyzing your question..."):
                response = get_chatbot_response(question, st.session_state.params, results)
                st.session_state.chat_history.append(("assistant", response))
            
            st.rerun()
        
        # Text input for custom questions
        user_input = st.text_input(
            "Type your question here:",
            placeholder="e.g., 'What if I reduce the number of participants to 15?' or 'Explain discounted payback period'",
            key="chat_input"
        )
        
        col1, col2 = st.columns([1, 4])
        with col1:
            send_button = st.button("Send 📤", type="primary")
        with col2:
            if st.button("Clear Chat 🗑️"):
                st.session_state.chat_history = []
                st.rerun()
        
        if send_button and user_input.strip():
            # Add user question to history
            st.session_state.chat_history.append(("user", user_input))
            
            # Get AI response
            with st.spinner("🤖 Thinking..."):
                response = get_chatbot_response(user_input, st.session_state.params, results)
                st.session_state.chat_history.append(("assistant", response))
            
            st.rerun()
        
        # Helpful suggestions
        if not st.session_state.chat_history:
            st.markdown("""
            **💡 Sample Questions to Get Started:**
            - "What does my {:.1f}% ROI mean in practical terms?"
            - "Why is my discounted payback longer than simple payback?"
            - "What cost of capital should I use for a startup?"
            - "How can I justify this investment to my CFO?"
            - "What if economic conditions change?"
            - "Should I run a pilot program first?"
            - "What are the biggest risks in my assumptions?"
            - "How do I measure success after implementation?"
            """.format(results['kpis']['roi']))
        
        # Context panel
        with st.expander("🔍 Current Analysis Context", expanded=False):
            col1, col2 = st.columns(2)
            
            with col1:
                st.markdown(f"""
                **Program Setup:**
                - {params['participants']} participants
                - {params['program_duration']} month duration
                - {format_currency(params['avg_salary'])} avg salary
                - {params['discount_rate']:.1f}% cost of capital
                """)
            
            with col2:
                st.markdown(f"""
                **Key Results:**
                - {results['kpis']['roi']:.1f}% ROI
                - {results['kpis']['payback_months']:.1f} month payback
                - {format_currency(results['kpis']['npv'])} NPV
                - {results['kpis']['benefit_cost_ratio']:.1f}:1 BCR
                """)
        
        # Pro tips
        st.info("""
        **💡 Pro Tips:**
        - Ask specific questions about your numbers for better insights
        - Use "What if..." questions to explore scenarios
        - Ask for explanations of financial terms you don't understand
        - Request industry comparisons and benchmarks
        - Get help interpreting results for different stakeholders
        """)
    
    # Export functionality
    st.divider()
    st.subheader("📄 Export Business Case")
    
    col1, col2, col3, col4, col5 = st.columns(5)
    
    with col1:
        if st.button("📊 Export Data (JSON)", type="secondary"):
            export_data = {
                'program': {k: v for k, v in st.session_state.params.items()},
                'results': results,
                'timestamp': datetime.now().isoformat()
            }
            st.download_button(
                label="Download JSON",
                data=json.dumps(export_data, indent=2),
                file_name=f"leadership_roi_analysis_{datetime.now().strftime('%Y%m%d_%H%M%S')}.json",
                mime="application/json"
            )
    
    with col2:
        if st.button("📄 Export Summary (CSV)", type="secondary"):
            dpb_value = f"{results['kpis']['discounted_payback_months']:.1f}" if results['kpis']['discounted_payback_months'] != float('inf') else "Never"
            
            summary_data = {
                'Metric': ['ROI (%)', 'Simple Payback (months)', 'Discounted Payback (months)', 'NPV ($)', 'Investment ($)', 'Annual Benefits ($)', 'Net Benefit ($)', 'IRR (%)', 'MIRR (%)'],
                'Value': [
                    f"{results['kpis']['roi']:.1f}",
                    f"{results['kpis']['payback_months']:.1f}",
                    dpb_value,
                    f"{results['kpis']['npv']:.0f}",
                    f"{results['costs']['total']:.0f}",
                    f"{results['benefits']['total_annual']:.0f}",
                    f"{results['kpis']['net_benefit']:.0f}",
                    f"{results['kpis']['irr']:.1f}" if results['kpis']['irr'] else "N/A",
                    f"{results['kpis']['mirr']:.1f}" if results['kpis']['mirr'] else "N/A"
                ]
            }
            df_summary = pd.DataFrame(summary_data)
            csv = df_summary.to_csv(index=False)
            st.download_button(
                label="Download CSV",
                data=csv,
                file_name=f"leadership_roi_summary_{datetime.now().strftime('%Y%m%d_%H%M%S')}.csv",
                mime="text/csv"
            )
    
    with col3:
        if st.button("📋 Export PDF Report", type="secondary"):
            try:
                pdf_buffer = generate_pdf_report(st.session_state.params, results)
                st.download_button(
                    label="Download PDF",
                    data=pdf_buffer,
                    file_name=f"leadership_roi_report_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pdf",
                    mime="application/pdf"
                )
            except Exception as e:
                st.error(f"Error generating PDF: {str(e)}")
                st.info("Install reportlab: pip install reportlab")
    
    with col4:
        if st.button("📊 Export PowerPoint", type="secondary"):
            try:
                ppt_buffer = generate_powerpoint_report(st.session_state.params, results)
                st.download_button(
                    label="Download PPTX",
                    data=ppt_buffer,
                    file_name=f"leadership_roi_presentation_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx",
                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                )
            except Exception as e:
                st.error(f"Error generating PowerPoint: {str(e)}")
                st.info("Install python-pptx: pip install python-pptx")
    
    with col5:
        recommendation = (
            "STRONG BUSINESS CASE - Proceed with implementation" if results['kpis']['roi'] >= 200 else
            "MODERATE BUSINESS CASE - Consider optimization" if results['kpis']['roi'] >= 100 else
            "WEAK BUSINESS CASE - Review assumptions and design"
        )
        st.success(f"**Recommendation:** {recommendation}")

if __name__ == "__main__":
    main()
